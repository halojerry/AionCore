"""build_deck.py — 端到端生成一份 4-5 页学术 PPTX（python-pptx）。

演示 light-slides 的完整产出链路：选 db06 主题 -> 行动式标题 -> 封面/
内容/结果/对比/References 五种版式 -> 每页 speaker notes -> 保存。
配色与字体全部取自 assets/themes.py，做到“审美统一、重点突出”。

运行：
    python build_deck.py                 # 用 academic 主题，出 academic_demo.pptx
    python build_deck.py --theme tech -o tech_demo.pptx

跑完后默认会打印每页标题（幽灵 deck 测试），方便检查论证连贯性。
本脚本不依赖任何外部数据，可直接自测运行。
"""
import sys
sys.stdout.reconfigure(encoding="utf-8")
import os
import argparse

# 让脚本能 import 同仓 assets/themes.py
HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, "..", "assets"))
from themes import get_theme  # noqa: E402

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

EMU = 914400


def C(hexstr):
    """hex(6位无#) -> RGBColor。"""
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def blank_slide(prs):
    """取空白版式（最后一个布局通常是 Blank），全部自由摆放。"""
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 \
        else prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)


def fill_bg(slide, prs, hexcolor):
    """铺满页背景色块（python-pptx 不能直接设页背景，用全幅矩形代替）。"""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = C(hexcolor)
    rect.line.fill.background()
    rect.shadow.inherit = False
    # 沉到底层
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)
    return rect


def add_band(slide, prs, hexcolor, top_in, h_in):
    """加一条横向色带（标题条 / 装饰条）。"""
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(top_in),
                                  prs.slide_width, Inches(h_in))
    band.fill.solid()
    band.fill.fore_color.rgb = C(hexcolor)
    band.line.fill.background()
    band.shadow.inherit = False
    return band


def _set_run_fonts(run, latin="Calibri", cjk="Microsoft YaHei"):
    """设拉丁 <a:latin> + 东亚 <a:ea>/<a:cs> 字体。只设 run.font.name 时中文回退默认黑体
    （中文优先技能 bug）——补 ea/cs 让中文用上主题中文字体，中英混排各用各的。"""
    from pptx.oxml.ns import qn
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", cjk)


def add_text(slide, x, y, w, h, text, size, color, *, bold=False,
             font="Calibri", cjk_font="Microsoft YaHei", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """加文本框（margin=0 便于与形状对齐）。text 可为 str 或 [(line, lvl)]。
    拉丁+CJK 字体都设，中文不回退默认黑体。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    lines = text if isinstance(text, list) else [(text, 0)]
    for i, item in enumerate(lines):
        line, lvl = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.level = lvl
        r = p.add_run()
        r.text = ("• " if lvl > 0 else "") + line
        r.font.size = Pt(size if lvl == 0 else max(14, size - 2))
        r.font.bold = bold and lvl == 0
        _set_run_fonts(r, latin=font, cjk=cjk_font)
        r.font.color.rgb = C(color)
    return tb


def add_accent_bar(slide, x, y, hexcolor, w_in=0.55, h_in=0.07):
    """标题左侧/下方的小色块强调（替代 AI 味下划线，做重复视觉母题）。"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(w_in), Inches(h_in))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C(hexcolor)
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar
def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def cover_slide(prs, t):
    """封面：深色块 + 大标题 + 副标题 + 作者单位日期。"""
    c, f = t["COLORS"], t["FONTS"]
    s = blank_slide(prs)
    fill_bg(s, prs, c["bg"])
    # 深色主视觉块占左 / 上
    add_band(s, prs, c["primary"], 0, prs.slide_height.inches * 0.62)
    add_text(s, 0.9, 1.4, 11.5, 1.4,
             "轻量化目标检测在边缘端的部署研究",
             40, c["bg"] if not t["dark"] else c["text"],
             bold=True, font=f["title"])
    add_accent_bar(s, 0.92, 2.95, c["accent"], w_in=1.4, h_in=0.09)
    add_text(s, 0.9, 3.15, 11.5, 0.8,
             "面向畜牧场景的实时检测与跟踪",
             20, c["secondary"] if t["dark"] else "FFFFFF", font=f["body"])
    add_text(s, 0.9, 5.1, 11.5, 1.2,
             [("李光  ·  西北农林科技大学  信息工程学院", 0),
              ("2026-06  ·  硕士学位论文答辩", 0)],
             16, c["text"], font=f["body"])
    set_notes(s, "开场 30s：自报家门，一句话点题——我们要让重模型在边缘设备实时跑起来。"
                 "不念标题，直接抛问题钩子。")
    return s


def content_slide(prs, t):
    """内容页：行动式标题 + 要点 bullet（左对齐、留白）。"""
    c, f = t["COLORS"], t["FONTS"]
    s = blank_slide(prs)
    fill_bg(s, prs, c["bg"])
    add_band(s, prs, c["surface"], 0, 1.25)
    add_text(s, 0.7, 0.32, 12, 0.8,
             "边缘端算力受限，重检测模型无法满足实时性",
             26, c["primary"], bold=True, font=f["title"])
    add_accent_bar(s, 0.72, 1.05, c["accent"])
    add_text(s, 0.9, 1.9, 11.5, 4.5,
             [("主流检测器参数量大，边缘 GPU 上帧率不足 15 FPS", 1),
              ("畜牧场景目标密集、遮挡多，小模型精度掉点明显", 1),
              ("现有压缩方法在精度与速度间难以兼顾", 1),
              ("目标：在 < 5M 参数下保持 mAP ≥ 0.85，达到 30 FPS", 1)],
             18, c["text"], font=f["body"])
    set_notes(s, "1min：讲清痛点三连——算力、密集遮挡、精度速度权衡，"
                 "最后落到我们的量化目标。必讲。")
    return s
def results_slide(prs, t):
    """结果页：左图右解读（图用原生柱状 chart，bullet 给 so-what）。"""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    c, f = t["COLORS"], t["FONTS"]
    s = blank_slide(prs)
    fill_bg(s, prs, c["bg"])
    add_band(s, prs, c["surface"], 0, 1.25)
    add_text(s, 0.7, 0.32, 12, 0.8,
             "本方法在精度持平下推理速度提升 2.3 倍",
             26, c["primary"], bold=True, font=f["title"])
    add_accent_bar(s, 0.72, 1.05, c["accent"])
    # 左：图
    cd = CategoryChartData()
    cd.categories = ["YOLO-base", "MobileNet", "Ours"]
    cd.add_series("FPS", (14, 22, 32))
    cd.add_series("mAP×100", (87, 79, 86))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                            Inches(0.7), Inches(1.7), Inches(6.6), Inches(4.6), cd)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True
    # 右：so-what 解读
    add_text(s, 7.6, 1.9, 5.2, 4.3,
             [("Ours 达 32 FPS，越过 30 FPS 实时线", 1),
              ("mAP 仅降 1 个点，精度基本无损", 1),
              ("参数量压到 4.6M，可部署 Jetson Nano", 1),
              ("↑2.3× 速度来自结构重参数 + 通道剪枝", 1)],
             18, c["text"], font=f["body"])
    set_notes(s, "1.5min：先读图——红柱是我们的；强调越过实时线这条关键结论，"
                 "再把右侧三点逐条对应。这是核心结果页，必讲且放慢。")
    return s


def compare_slide(prs, t):
    """对比页：三列方法对比表，突出本方法列。"""
    c, f = t["COLORS"], t["FONTS"]
    s = blank_slide(prs)
    fill_bg(s, prs, c["bg"])
    add_band(s, prs, c["surface"], 0, 1.25)
    add_text(s, 0.7, 0.32, 12, 0.8,
             "相比两类基线，本方法同时赢得速度、体积与可部署性",
             24, c["primary"], bold=True, font=f["title"])
    add_accent_bar(s, 0.72, 1.05, c["accent"])
    rows, cols = 5, 4
    tb = s.shapes.add_table(rows, cols, Inches(0.9), Inches(1.8),
                            Inches(11.5), Inches(4.0)).table
    headers = ["维度", "YOLO-base", "MobileNet", "Ours"]
    data = [
        ["参数量 (M)", "37.2", "6.8", "4.6"],
        ["速度 (FPS)", "14", "22", "32"],
        ["mAP", "0.87", "0.79", "0.86"],
        ["边缘可部署", "否", "勉强", "是"],
    ]
    for j, htext in enumerate(headers):
        cell = tb.cell(0, j)
        cell.text = htext
        cell.fill.solid()
        cell.fill.fore_color.rgb = C(c["primary"])
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(15)
        para.font.bold = True
        para.font.color.rgb = C("FFFFFF")
        para.font.name = f["body"]
    for i, row in enumerate(data, 1):
        for j, val in enumerate(row):
            cell = tb.cell(i, j)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14)
            para.font.name = f["body"]
            highlight = (j == 3)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C(c["accent"]) if highlight else C(c["surface"])
            para.font.color.rgb = C("FFFFFF") if highlight else C(c["text"])
            para.font.bold = highlight
    set_notes(s, "1min：表格只强调最后一列 Ours（高亮色）。逐行对比，"
                 "落点在『边缘可部署=是』。可略讲中间数值。")
    return s


def references_slide(prs, t):
    """References 页：编号文献列表（学术 deck 收尾前的引用页）。"""
    c, f = t["COLORS"], t["FONTS"]
    s = blank_slide(prs)
    fill_bg(s, prs, c["bg"])
    add_band(s, prs, c["surface"], 0, 1.25)
    add_text(s, 0.7, 0.32, 12, 0.8, "References",
             26, c["primary"], bold=True, font=f["en"])
    add_accent_bar(s, 0.72, 1.05, c["accent"])
    refs = [
        "[1] Redmon J, Farhadi A. YOLOv3: An Incremental Improvement. arXiv:1804.02767, 2018.",
        "[2] Howard A, et al. MobileNetV3. ICCV, 2019.",
        "[3] Ding X, et al. RepVGG: Making VGG-style ConvNets Great Again. CVPR, 2021.",
        "[4] He Y, et al. Channel Pruning for Accelerating Networks. ICCV, 2017.",
    ]
    add_text(s, 0.9, 1.9, 11.8, 4.4, [(r, 0) for r in refs],
             14, c["muted"], font=f["en"])
    set_notes(s, "15s：引用页一般不展开念，留作存档与提问时回查。"
                 "占位文献仅为演示，真实交付须逐条核实 DOI（CONVENTIONS §4）。")
    return s
def conclusions_slide(prs, t):
    """Conclusions 页：学术 deck 的正确收尾（不写 'Thank You/Q&A' 占整页——
    最后一页应是带走的结论。对齐 SKILL「Conclusions 收尾」规则）。"""
    c, f = t["COLORS"], t["FONTS"]
    s = blank_slide(prs)
    fill_bg(s, prs, c["bg"])
    add_band(s, prs, c["surface"], 0, 1.25)
    add_text(s, 0.7, 0.32, 12, 0.8, "Conclusions",
             26, c["primary"], bold=True, font=f["en"])
    add_accent_bar(s, 0.72, 1.05, c["accent"])
    points = [
        ("我们提出 X，在 3 数据集上较最强 baseline 提升 4.2%（p<0.01）", 0),
        ("消融证实各模块均有正贡献，增益不来自参数量", 0),
        ("局限：稀有类样本不足时提前量增益下降；未来补主动采集", 0),
        ("Takeaway: 对比学习 + 领域先验是发情早期识别的有效路径", 0),
    ]
    add_text(s, 0.9, 1.9, 11.6, 4.2, points, 20, c["text"],
             font=f["en"], cjk_font=f.get("cjk", "Microsoft YaHei"))
    set_notes(s, "40s：用结论页收尾而非 Thank You。每条对应一个贡献/局限，"
                 "最后一句给评委能带走的 takeaway。")
    return s
def build(theme_name="academic", out="academic_demo.pptx"):
    t = get_theme(theme_name)
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 宽屏
    prs.slide_height = Inches(7.5)
    # 学术顺序：封面→内容→结果→对比→结论(收尾带走点)→引用(存档回查)
    builders = [cover_slide, content_slide, results_slide,
                compare_slide, conclusions_slide, references_slide]
    slides = [b(prs, t) for b in builders]
    prs.save(out)
    return prs, out, t


def ghost_deck_test(prs):
    """幽灵 deck 测试：抽出每页首个非空文本作为标题，连起来看论证是否连贯。"""
    titles = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().split("\n")[0]
                break
        titles.append(f"  {i}. {title}")
    return titles


def main():
    ap = argparse.ArgumentParser(description="端到端生成 5 页学术 PPTX")
    ap.add_argument("--theme", default="academic")
    ap.add_argument("-o", "--out", default=None)
    args = ap.parse_args()
    out = args.out or f"{args.theme}_demo.pptx"
    prs, out, t = build(args.theme, out)
    size_kb = os.path.getsize(out) / 1024
    print(f"已生成：{out}（主题 {t['name']} / {t['label']}，"
          f"{len(prs.slides._sldIdLst)} 页，{size_kb:.0f} KB）")
    print("幽灵 deck 测试（连读标题应能讲完论证）：")
    for line in ghost_deck_test(prs):
        print(line)


if __name__ == "__main__":
    main()
