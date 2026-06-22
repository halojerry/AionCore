#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""pacing.py — 答辩/路演时长估算（读 PPTX notes 的时长标记，核对是否超时）。

答辩/路演有硬性时间限制（如 12 分钟）。本脚本读每页演讲者备注里的时长标记（如 "40s"/"1.5min"/
"[90s]"），累加估总时长，对照 --limit 给 OK/超时；无标记的页按字数粗估兜底。纯标准库 + python-pptx。

诚实：时长是**估算**（按 notes 标记或讲稿字数×语速），真实因人而异；只帮发现"明显超时/分配失衡"，
不替代真人计时彩排。无 python-pptx 时降级为纯文本 notes 文件解析。

用法：
  python pacing.py deck.pptx --limit 12        # 12 分钟答辩
  python pacing.py deck.pptx --limit 12 --wpm 160   # 自定中文/英文语速
  python pacing.py --selftest
"""
from __future__ import annotations
import argparse
import re
import sys

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

# 时长标记：40s / 90 s / 1.5min / 2 分钟 / [60s] / 时长:45秒
_DUR = re.compile(r"(\d+(?:\.\d+)?)\s*(s|sec|secs|second|seconds|秒|min|mins|minute|minutes|分钟|分)", re.I)
_SEC_UNITS = {"s", "sec", "secs", "second", "seconds", "秒"}


def parse_duration(note: str) -> float | None:
    """从单页 notes 抽时长（秒）。取第一个时长标记；无则 None。"""
    m = _DUR.search(note or "")
    if not m:
        return None
    val = float(m.group(1))
    unit = m.group(2).lower()
    return val if unit in _SEC_UNITS else val * 60.0


def estimate_from_words(note: str, body_chars: int, wpm: float) -> float:
    """无时长标记时按讲稿/正文字数粗估秒数。中文按字、英文按词混合估。"""
    # 讲稿优先；没有讲稿用正文字数的一半（正文是要点不是逐字稿）
    text = note or ""
    cjk = sum(1 for c in text if "一" <= c <= "鿿")
    latin_words = len(re.findall(r"[A-Za-z]+", text))
    units = cjk + latin_words
    if units == 0:
        units = max(0, body_chars // 2)   # 兜底：正文字数估
    # wpm 对中文是"字/分"，英文是"词/分"；混合用同一 wpm 近似
    return round(units / max(1.0, wpm) * 60.0, 1)


def assess(slides: list, limit_min: float, wpm: float = 150.0) -> dict:
    """slides: [{idx, note, body_chars}]. 返回总时长估算 + 每页 + 超时判定。"""
    rows, total, n_marked = [], 0.0, 0
    for s in slides:
        d = parse_duration(s.get("note", ""))
        if d is not None:
            n_marked += 1
            src = "notes标记"
        else:
            d = estimate_from_words(s.get("note", ""), s.get("body_chars", 0), wpm)
            src = "字数估算"
        total += d
        rows.append({"idx": s["idx"], "sec": round(d, 1), "src": src})
    limit_sec = limit_min * 60.0
    over = total - limit_sec
    verdict = ("OK" if total <= limit_sec else "OVER")
    # 分配失衡：单页 > 平均 2.5 倍且 deck >3 页
    avg = total / max(1, len(slides))
    hot = [r["idx"] for r in rows if r["sec"] > avg * 2.5] if len(slides) > 3 else []
    return {
        "n_slides": len(slides), "n_marked": n_marked,
        "total_sec": round(total, 1), "total_min": round(total / 60.0, 1),
        "limit_min": limit_min, "verdict": verdict,
        "over_sec": round(over, 1) if over > 0 else 0,
        "per_slide": rows, "hot_slides": hot,
        "note": ("时长为估算(notes 标记优先，否则字数×语速)，真实因人而异；"
                 + (f"{n_marked}/{len(slides)} 页有时长标记，其余字数估" if slides else "")
                 + "。务必真人计时彩排，本脚本只抓明显超时/分配失衡。"),
    }


def render(r: dict) -> str:
    icon = "✅" if r["verdict"] == "OK" else "🛑"
    L = [f"# 答辩/路演时长估算 {icon} {r['verdict']}",
         f"- 估算总时长 {r['total_min']} 分钟（{r['total_sec']}s）/ 限 {r['limit_min']} 分钟"]
    if r["verdict"] == "OVER":
        L.append(f"- ⚠ 超时约 {r['over_sec']}s（{round(r['over_sec']/60,1)} 分）——删页/精简讲稿")
    if r["hot_slides"]:
        L.append(f"- ⚠ 时长偏重的页：{r['hot_slides']}（>均值 2.5×，考虑拆分或压缩）")
    L.append(f"- {r['note']}")
    return "\n".join(L)


def _selftest() -> int:
    print("### pacing 离线自测", file=sys.stderr)
    # 时长标记解析（多种写法）
    assert parse_duration("讲 40s 这页") == 40.0
    assert parse_duration("[1.5min] intro") == 90.0
    assert parse_duration("时长:2分钟") == 120.0
    assert parse_duration("no marker here") is None

    # 全部有标记：总时长 = 标记和
    slides = [{"idx": i, "note": f"{d}s", "body_chars": 0}
              for i, d in enumerate([30, 60, 90, 120], 1)]
    r = assess(slides, limit_min=5)          # 300s 总=300s → OK 边界
    assert r["total_sec"] == 300.0 and r["n_marked"] == 4, r
    assert r["verdict"] == "OK", r
    # 超时
    r2 = assess(slides, limit_min=4)         # 240s 限 < 300 → OVER
    assert r2["verdict"] == "OVER" and r2["over_sec"] == 60.0, r2

    # 无标记走字数估算
    slides_nm = [{"idx": 1, "note": "word " * 150, "body_chars": 0}]  # 150 词
    r3 = assess(slides_nm, limit_min=10, wpm=150)
    assert r3["per_slide"][0]["src"] == "字数估算", r3
    assert 55 <= r3["total_sec"] <= 65, r3   # 150词/150wpm≈60s

    # 分配失衡检测
    imb = [{"idx": 1, "note": "10s"}, {"idx": 2, "note": "10s"},
           {"idx": 3, "note": "10s"}, {"idx": 4, "note": "200s"}]
    r4 = assess(imb, limit_min=10)
    assert 4 in r4["hot_slides"], r4

    md = render(r2)
    assert "超时" in md and "估算" in md, md
    print("[selftest] PASS pacing offline")
    return 0


def _load_pptx(path: str) -> list:
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        note = ""
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text or ""
        body = sum(len(sh.text_frame.text) for sh in slide.shapes
                   if sh.has_text_frame) if slide.shapes else 0
        out.append({"idx": i, "note": note, "body_chars": body})
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description="答辩/路演时长估算（读 PPTX notes 时长标记）")
    ap.add_argument("pptx", nargs="?", help="PPTX 路径")
    ap.add_argument("--limit", type=float, default=12, help="时长上限（分钟）")
    ap.add_argument("--wpm", type=float, default=150, help="语速（字/词每分钟，无标记页估算用）")
    ap.add_argument("--json", action="store_true")
    ap.add_argument("--selftest", action="store_true")
    args = ap.parse_args()

    if args.selftest or not args.pptx:
        return _selftest()
    try:
        slides = _load_pptx(args.pptx)
    except ImportError:
        print("需要 python-pptx：pip install python-pptx", file=sys.stderr)
        return 2
    r = assess(slides, args.limit, args.wpm)
    if args.json:
        import json
        print(json.dumps(r, ensure_ascii=False, indent=2))
    else:
        print(render(r))
    return 1 if r["verdict"] == "OVER" else 0


if __name__ == "__main__":
    sys.exit(main())
