#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""pptx_eval.py — PPT 可量化质量评测（借 PPTAgent/PPTEval 思路，落成 Light 能算的指标）。

thumbnail.py 出缩略图给人肉眼查；本脚本把"幽灵测试/视觉 QA"升级为**可打分、可回归**的
三维评测，对照 SKILL「具体尺度」里的硬数值，逐页算分并给扣分理由：

  - content 内容密度：每页 bullet 数、字数、是否纯文字页/空页（对齐 SKILL：每页≤7 bullet、
    每页须有视觉元素、禁空页）
  - design 设计一致：全 deck 配色数（主+1~2 辅+1 强调，≤4~5）、字号档位是否在区间
    （标题 36–44 / 节标题 20–24 / 正文 14–16 / 注释 10–12 pt）
  - coherence 连贯性：版式重复度（禁每页同一版式）、母题一致

⚠ 诚实声明：
- 这是**结构性可计算指标**，不评"美不美/内容对不对"——颜色搭配品味、叙事质量仍须人工。
- 分数是**启发式扣分制**（满分 10，违反硬尺度按条扣），阈值对齐 SKILL 经验尺度、可调，非真值。
- 纯 python 读 pptx 几何/run 字号/填充色，不渲染像素；与 thumbnail.py 互补（它给人看、本脚本给分）。

用法：
  python pptx_eval.py deck.pptx
  python pptx_eval.py deck.pptx --json
  python pptx_eval.py --selftest
"""
import sys
sys.stdout.reconfigure(encoding="utf-8")
import argparse
import json

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
except ImportError:
    Presentation = None

EMU_PER_INCH = 914400

# 评测阈值（对齐 SKILL「具体尺度」，经验值、可调，非真值）
MAX_BULLETS_PER_SLIDE = 7      # 每页要点上限
MAX_WORDS_PER_SLIDE = 110      # 每页字数软上限（中英混合粗估）
MAX_DECK_COLORS = 5            # 全 deck 主色族上限（主+1~2辅+1强调）
FONT_BANDS = {                 # 字号合理区间（pt），落区间外记疑似
    "title": (28, 48), "body": (12, 20), "note": (9, 13),
}


def _iter_text_runs(shape):
    """产出 (text, font_size_pt_or_None)。仅文本框/占位符。"""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            sz = None
            if run.font.size is not None:
                sz = run.font.size.pt
            yield (run.text or ""), sz


def _shape_fill_rgb(shape):
    """尽力取形状填充色 RGB 十六进制（取不到返回 None，不臆造）。"""
    try:
        fill = shape.fill
        if fill.type is not None and hasattr(fill, "fore_color"):
            rgb = fill.fore_color.rgb
            if rgb is not None:
                return str(rgb)
    except Exception:  # noqa: 颜色取不到（主题色/继承）很常见，不报错
        return None
    return None


def analyze_slide(slide) -> dict:
    """抽单页结构特征：bullet 数、字数、视觉元素数、字号集合、填充色集合、版式名。"""
    bullets = 0
    words = 0
    font_sizes = []
    colors = set()
    visual_shapes = 0   # 非文本的图片/形状/图表（判"禁纯文字页"）
    text_shapes = 0
    for shape in slide.shapes:
        is_text = getattr(shape, "has_text_frame", False) and (shape.text_frame.text or "").strip()
        if is_text:
            text_shapes += 1
            for para in shape.text_frame.paragraphs:
                ptext = "".join(r.text or "" for r in para.runs) or para.text or ""
                if ptext.strip():
                    bullets += 1
                    # 字数：英文按空格、中文按字符粗估
                    words += len(ptext.split()) + sum(1 for c in ptext if "一" <= c <= "鿿")
            for _, sz in _iter_text_runs(shape):
                if sz:
                    font_sizes.append(sz)
        else:
            # 图片/图表/自选图形等视觉元素
            stype = str(getattr(shape, "shape_type", "") or "")
            if "PICTURE" in stype or "CHART" in stype or "AUTO_SHAPE" in stype or "GROUP" in stype:
                visual_shapes += 1
        c = _shape_fill_rgb(shape)
        if c:
            colors.add(c)
    layout_name = ""
    try:
        layout_name = slide.slide_layout.name or ""
    except Exception:
        pass
    return {"bullets": bullets, "words": words, "font_sizes": font_sizes,
            "colors": colors, "visual_shapes": visual_shapes,
            "text_shapes": text_shapes, "layout": layout_name}


def evaluate(path: str) -> dict:
    if Presentation is None:
        raise RuntimeError("需要 python-pptx：pip install python-pptx")
    prs = Presentation(path)
    slides = [analyze_slide(s) for s in prs.slides]
    n = len(slides)
    if n == 0:
        return {"error": "空 deck（0 页）", "n_slides": 0}

    deck_colors = set()
    layouts = []
    content_issues, design_issues, coherence_issues = [], [], []

    for i, s in enumerate(slides, 1):
        deck_colors |= s["colors"]
        layouts.append(s["layout"])
        # content 维度
        if s["bullets"] > MAX_BULLETS_PER_SLIDE:
            content_issues.append(f"P{i}: {s['bullets']} 个要点 > {MAX_BULLETS_PER_SLIDE}（拆页或精简）")
        if s["words"] > MAX_WORDS_PER_SLIDE:
            content_issues.append(f"P{i}: ~{s['words']} 字 > {MAX_WORDS_PER_SLIDE}（信息过载）")
        if s["text_shapes"] > 0 and s["visual_shapes"] == 0 and i > 1:
            content_issues.append(f"P{i}: 疑似纯文字页（无图片/图表/形状，SKILL 禁纯文字页）")
        if s["bullets"] == 0 and s["visual_shapes"] == 0:
            content_issues.append(f"P{i}: 疑似空页")
        # design 维度：字号落区间外
        for sz in s["font_sizes"]:
            in_band = any(lo <= sz <= hi for lo, hi in FONT_BANDS.values())
            if not in_band:
                design_issues.append(f"P{i}: 字号 {sz}pt 不在常用区间（标题28-48/正文12-20/注释9-13）")
                break

    # design：全 deck 配色数
    if len(deck_colors) > MAX_DECK_COLORS:
        design_issues.append(f"全 deck 取到 {len(deck_colors)} 种填充色 > {MAX_DECK_COLORS}（配色应主+1~2辅+1强调）")
    # coherence：版式重复度（禁每页同一版式；但封面/结论可不同）
    nonempty_layouts = [l for l in layouts if l]
    if nonempty_layouts:
        most = max(set(nonempty_layouts), key=nonempty_layouts.count)
        rep = nonempty_layouts.count(most) / len(nonempty_layouts)
        if rep > 0.8 and len(nonempty_layouts) >= 5:
            coherence_issues.append(f"{rep:.0%} 页用同一版式 '{most}'（SKILL 禁每页同版式，缺层次）")
        uniq = len(set(nonempty_layouts))
        if uniq == 1 and len(nonempty_layouts) >= 4:
            coherence_issues.append("全 deck 仅 1 种版式（封面/内容/结论应有别）")

    # 扣分制：每维满分 10，每条 issue 扣分（封顶扣到 0）
    def score(issues, per=2.0):
        return max(0.0, round(10 - per * len(issues), 1))
    content_s = score(content_issues)
    design_s = score(design_issues)
    coherence_s = score(coherence_issues, per=3.0)
    overall = round((content_s + design_s + coherence_s) / 3, 1)

    return {
        "n_slides": n,
        "scores": {"content": content_s, "design": design_s,
                   "coherence": coherence_s, "overall": overall},
        "deck_color_count": len(deck_colors),
        "issues": {"content": content_issues, "design": design_issues,
                   "coherence": coherence_issues},
        "note": ("可计算的结构性指标(扣分制，阈值对齐 SKILL 经验尺度、可调)；"
                 "不评美感/内容正确性——配色品味、叙事质量、图表是否真数据仍须人工 + thumbnail.py 肉眼复核。"),
    }


def to_markdown(rep: dict) -> str:
    if rep.get("error"):
        return f"# PPT 评测\n\n{rep['error']}"
    s = rep["scores"]
    lines = [f"# PPT 量化评测 — {rep['n_slides']} 页\n",
             f"| 维度 | 分(/10) |", "|---|---|",
             f"| 内容密度 content | {s['content']} |",
             f"| 设计一致 design | {s['design']} |",
             f"| 连贯性 coherence | {s['coherence']} |",
             f"| **综合 overall** | **{s['overall']}** |", ""]
    for dim, label in (("content", "内容"), ("design", "设计"), ("coherence", "连贯性")):
        iss = rep["issues"][dim]
        if iss:
            lines.append(f"**{label}问题（{len(iss)}）**：")
            lines += [f"- {x}" for x in iss]
            lines.append("")
    lines.append(f"> {rep['note']}")
    return "\n".join(lines)


def _build_demo_pptx(path, bad=False):
    """构造合成 pptx 供 selftest（不联网）。bad=True 故意造内容超载+单一版式。"""
    from pptx import Presentation as P
    from pptx.util import Inches, Pt as _Pt
    prs = P()
    blank = prs.slide_layouts[6]
    titled = prs.slide_layouts[5]
    if bad:
        # 5 页都用同一空白版式、每页塞 9 个要点、无视觉元素 → 触发内容+连贯性多维扣分
        for _ in range(5):
            sl = prs.slides.add_slide(blank)
            tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(6))
            tf = tb.text_frame
            for k in range(9):
                p = tf.add_paragraph() if k else tf.paragraphs[0]
                p.text = f"这是第 {k+1} 个要点，内容比较长用来凑足够多的字数触发信息过载阈值的检测逻辑啊"
                p.runs[0].font.size = _Pt(15)
    else:
        # 好 deck：封面(标题版式) + 2 内容页各 3 要点 + 一个形状(视觉元素)
        s0 = prs.slides.add_slide(titled)
        s0.shapes.title.text = "演示标题"
        for _ in range(2):
            sl = prs.slides.add_slide(blank)
            tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(4))
            tf = tb.text_frame
            for k in range(3):
                p = tf.add_paragraph() if k else tf.paragraphs[0]
                p.text = f"要点 {k+1}"
                p.runs[0].font.size = _Pt(16)
            # 加个自选图形当视觉元素
            from pptx.enum.shapes import MSO_SHAPE
            sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(1), Inches(2), Inches(2))
    prs.save(path)


def _selftest() -> int:
    print("### pptx_eval 离线自测", file=sys.stderr)
    if Presentation is None:
        print("[selftest] SKIP：未装 python-pptx", file=sys.stderr)
        return 0
    import tempfile, os
    d = tempfile.mkdtemp()
    good_p, bad_p = os.path.join(d, "good.pptx"), os.path.join(d, "bad.pptx")
    _build_demo_pptx(good_p, bad=False)
    _build_demo_pptx(bad_p, bad=True)

    good = evaluate(good_p)
    bad = evaluate(bad_p)
    print(f"[good] overall={good['scores']['overall']} content={good['scores']['content']}")
    print(f"[bad ] overall={bad['scores']['overall']} issues={sum(len(v) for v in bad['issues'].values())}")

    # 好 deck 综合分应明显高于坏 deck
    assert good["scores"]["overall"] > bad["scores"]["overall"], (good["scores"], bad["scores"])
    # 坏 deck 应触发 content（9>7 要点 + 信息过载）
    assert any("要点" in x for x in bad["issues"]["content"]), bad["issues"]["content"]
    # 坏 deck 3 页同空白版式 → coherence 问题
    assert bad["issues"]["coherence"], bad["issues"]
    # markdown 含诚实声明
    md = to_markdown(good)
    assert "人工" in md and "thumbnail" in md, md
    # 空 deck 不崩
    empty_p = os.path.join(d, "empty.pptx")
    Presentation().save(empty_p)
    assert evaluate(empty_p).get("n_slides") == 0

    import shutil
    shutil.rmtree(d, ignore_errors=True)
    print("[selftest] PASS pptx_eval offline")
    return 0


def main() -> int:
    ap = argparse.ArgumentParser(description="PPT 可量化质量评测（内容/设计/连贯性三维扣分制）")
    ap.add_argument("pptx", nargs="?", help="待评测 .pptx 路径")
    ap.add_argument("--json", action="store_true", help="输出 JSON")
    ap.add_argument("--selftest", action="store_true", help="run offline synthetic self-test")
    args = ap.parse_args()

    if args.selftest or not args.pptx:
        return _selftest()
    rep = evaluate(args.pptx)
    print(json.dumps(rep, ensure_ascii=False, indent=2, default=list) if args.json
          else to_markdown(rep))
    return 0


if __name__ == "__main__":
    sys.exit(main())

