"""thumbnail.py — 把 .pptx 渲染成缩略图网格，做整套 deck 的快速视觉 QA。

两条渲染路径，自动择优：
  1) 若系统有 LibreOffice（soffice），先转 PDF 再用 pdftoppm/PyMuPDF 出真实位图（最准）。
  2) 无 soffice 时（本环境默认）走纯 python 回退：用 python-pptx 读每页形状的
     位置/尺寸/填充色/文字，用 PIL 画出版式示意图（boxes + 文字裁剪），
     足以抓重叠/溢出/空页/对齐错位/版式重复等结构性问题。

用法：
    python thumbnail.py deck.pptx                 # 出 deck_thumbnails.png
    python thumbnail.py deck.pptx -o grid.png --cols 3 --scale 1.5
    python thumbnail.py deck.pptx --engine pillow # 强制纯 python 回退

依赖：python-pptx, Pillow（均已在环境内）。soffice/PyMuPDF 可选。
注意：纯 python 回退是“版式示意”，不渲染渐变/图片像素/字体精排，
      用于结构 QA；要像素级保真请装 LibreOffice 后用 --engine soffice。
"""
import sys
sys.stdout.reconfigure(encoding="utf-8")
import os
import math
import shutil
import argparse
import subprocess
import tempfile

from PIL import Image, ImageDraw, ImageFont

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    Presentation = None

EMU_PER_INCH = 914400


def _find_soffice():
    """返回可用的 soffice 路径或 None。"""
    for cand in ("soffice", "libreoffice",
                 r"C:\Program Files\LibreOffice\program\soffice.exe",
                 r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"):
        if shutil.which(cand) or os.path.exists(cand):
            return cand
    return None


def _load_font(size):
    """尽量加载一个可读字体；失败回退 PIL 默认位图字体。"""
    for name in ("msyh.ttc", "msyhl.ttc", "simhei.ttf", "arial.ttf",
                 "DejaVuSans.ttf"):
        try:
            return ImageFont.truetype(name, size)
        except OSError:
            continue
    return ImageFont.load_default()


def _hex_to_rgb(h, default=(127, 127, 127)):
    if not h:
        return default
    h = str(h).lstrip("#")
    if len(h) != 6:
        return default
    try:
        return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return default


def _shape_fill_hex(shape):
    """尽力取形状填充色 hex（取不到返回 None）。"""
    try:
        fill = shape.fill
        if fill.type is not None and fill.fore_color and fill.fore_color.type is not None:
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None
def render_slides_pillow(pptx_path, px_w=480):
    """纯 python 回退：用 python-pptx 几何 + PIL 画每页版式示意图。
    返回 [PIL.Image, ...]，每页一张。"""
    if Presentation is None:
        raise RuntimeError("需要 python-pptx：pip install python-pptx")
    prs = Presentation(pptx_path)
    sw = prs.slide_width or Emu(int(13.333 * EMU_PER_INCH))
    sh = prs.slide_height or Emu(int(7.5 * EMU_PER_INCH))
    px_h = int(px_w * sh / sw)
    sx, sy = px_w / sw, px_h / sh
    imgs = []
    for idx, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (px_w, px_h), (255, 255, 255))
        d = ImageDraw.Draw(img)
        # 页背景：若首形状是铺满的色块，用其填充色当背景
        for shape in slide.shapes:
            try:
                if (shape.left in (0, None) and shape.top in (0, None)
                        and shape.width and shape.width >= sw * 0.98):
                    bg = _shape_fill_hex(shape)
                    if bg:
                        d.rectangle([0, 0, px_w, px_h], fill=_hex_to_rgb(bg))
                        break
            except Exception:
                pass
        for shape in slide.shapes:
            try:
                left = int((shape.left or 0) * sx)
                top = int((shape.top or 0) * sy)
                right = left + int((shape.width or 0) * sx)
                bot = top + int((shape.height or 0) * sy)
            except Exception:
                continue
            fill = _shape_fill_hex(shape)
            has_text = shape.has_text_frame and shape.text_frame.text.strip()
            if fill and (right - left) > 2 and (bot - top) > 2:
                d.rectangle([left, top, right, bot], fill=_hex_to_rgb(fill))
            elif shape.has_text_frame is False:
                # 图片/表格等：画占位描边框
                d.rectangle([left, top, right, bot], outline=(170, 170, 170))
            if has_text:
                txt = shape.text_frame.text.strip().replace("\n", " ")
                fnt = _load_font(max(9, min(16, (bot - top) // 3 or 11)))
                tcol = (30, 30, 30)
                # 文字色尽量取第一段第一 run
                try:
                    r0 = shape.text_frame.paragraphs[0].runs[0]
                    if r0.font.color and r0.font.color.type is not None:
                        tcol = _hex_to_rgb(str(r0.font.color.rgb), tcol)
                except Exception:
                    pass
                # 裁剪到框宽
                maxc = max(4, (right - left) // 7)
                d.text((left + 3, top + 2), txt[:maxc], fill=tcol, font=fnt)
        # 页码角标
        d.text((px_w - 22, px_h - 16), str(idx), fill=(150, 150, 150),
               font=_load_font(11))
        imgs.append(img)
    return imgs


def render_slides_soffice(pptx_path, px_w=480):
    """LibreOffice 路径：pptx -> pdf -> 位图。需 soffice + (PyMuPDF 或 pdftoppm)。"""
    soffice = _find_soffice()
    if not soffice:
        raise RuntimeError("未找到 soffice")
    tmp = tempfile.mkdtemp(prefix="thumb_")
    subprocess.run([soffice, "--headless", "--convert-to", "pdf",
                    "--outdir", tmp, pptx_path],
                   check=True, capture_output=True, timeout=180)
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    pdf = os.path.join(tmp, base + ".pdf")
    imgs = []
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(pdf)
        for page in doc:
            zoom = px_w / page.rect.width
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            imgs.append(Image.frombytes("RGB", [pix.width, pix.height], pix.samples))
        doc.close()
    except ImportError:
        if not shutil.which("pdftoppm"):
            raise RuntimeError("有 soffice 但缺 PyMuPDF 与 pdftoppm，无法栅格化 PDF")
        subprocess.run(["pdftoppm", "-png", "-r", "96", pdf,
                        os.path.join(tmp, "p")], check=True)
        for f in sorted(os.listdir(tmp)):
            if f.startswith("p") and f.endswith(".png"):
                imgs.append(Image.open(os.path.join(tmp, f)).convert("RGB"))
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    return imgs


def make_grid(imgs, cols=4, pad=12, label=True):
    """把每页缩略图拼成网格大图，逐格描边 + 标页码。"""
    if not imgs:
        raise RuntimeError("没有可渲染的幻灯片")
    cols = max(1, min(cols, len(imgs)))
    rows = math.ceil(len(imgs) / cols)
    cw = max(i.width for i in imgs)
    ch = max(i.height for i in imgs)
    lab_h = 18 if label else 0
    W = cols * cw + (cols + 1) * pad
    H = rows * (ch + lab_h) + (rows + 1) * pad
    grid = Image.new("RGB", (W, H), (245, 245, 247))
    d = ImageDraw.Draw(grid)
    fnt = _load_font(12)
    for i, im in enumerate(imgs):
        r, c = divmod(i, cols)
        x = pad + c * (cw + pad)
        y = pad + r * (ch + lab_h + pad)
        grid.paste(im, (x, y))
        d.rectangle([x, y, x + im.width, y + im.height], outline=(200, 200, 200))
        if label:
            d.text((x + 2, y + im.height + 2), f"Slide {i + 1}",
                   fill=(90, 90, 90), font=fnt)
    return grid



def _selftest() -> int:
    if Presentation is None:
        print("[selftest] SKIP thumbnail: python-pptx not installed")
        return 0
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for idx, color in enumerate(("4472C4", "70AD47", "ED7D31"), 1):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_shape(1, Inches(0.7), Inches(0.7), Inches(4.5), Inches(1.2))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string(color)
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(6), Inches(1))
        tx.text = f"Slide {idx}"
    with tempfile.TemporaryDirectory(prefix="light_thumb_") as tmp:
        pptx_path = os.path.join(tmp, "demo.pptx")
        out = os.path.join(tmp, "grid.png")
        prs.save(pptx_path)
        imgs = render_slides_pillow(pptx_path, px_w=240)
        assert len(imgs) == 3, len(imgs)
        grid = make_grid(imgs, cols=2)
        grid.save(out)
        assert os.path.exists(out) and os.path.getsize(out) > 0, out
    print("[selftest] PASS thumbnail")
    return 0


def main():
    ap = argparse.ArgumentParser(description="把 pptx 渲染成缩略图网格做视觉 QA")
    ap.add_argument("pptx", nargs="?")
    ap.add_argument("-o", "--out", help="输出 png（默认 <deck>_thumbnails.png）")
    ap.add_argument("--cols", type=int, default=4)
    ap.add_argument("--width", type=int, default=480, help="单页像素宽")
    ap.add_argument("--scale", type=float, default=1.0, help="单页宽缩放系数")
    ap.add_argument("--engine", choices=["auto", "soffice", "pillow"],
                    default="auto")
    ap.add_argument("--selftest", action="store_true", help="run synthetic pptx thumbnail self-test")
    args = ap.parse_args()

    if args.selftest:
        raise SystemExit(_selftest())
    if not args.pptx:
        ap.error("需要提供 pptx 路径（或使用 --selftest）")

    if not os.path.exists(args.pptx):
        sys.exit(f"找不到文件：{args.pptx}")
    px_w = int(args.width * args.scale)

    engine = args.engine
    imgs = None
    if engine in ("auto", "soffice") and _find_soffice():
        try:
            imgs = render_slides_soffice(args.pptx, px_w)
            engine = "soffice"
        except Exception as e:
            print(f"[thumbnail] soffice 路径失败（{e}），回退 pillow")
            imgs = None
    if imgs is None:
        if engine == "soffice":
            sys.exit("指定 soffice 引擎但不可用；装 LibreOffice 或改 --engine pillow")
        imgs = render_slides_pillow(args.pptx, px_w)
        engine = "pillow"

    grid = make_grid(imgs, cols=args.cols)
    out = args.out or (os.path.splitext(args.pptx)[0] + "_thumbnails.png")
    grid.save(out)
    note = "（版式示意，非像素级）" if engine == "pillow" else "（LibreOffice 真实渲染）"
    print(f"引擎={engine} {note}  共 {len(imgs)} 页 -> {out}  尺寸 {grid.size}")


if __name__ == "__main__":
    main()
